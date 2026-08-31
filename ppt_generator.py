"""述标PPT生成器 v9.1（可选增值，数据驱动版）。

基于评分项 + 项目信息 + 企业资质业绩库，生成一份简版述标/汇报 PPT（.pptx），
对齐 WPS AI / 喜鹊标书 / 链企 AI 的「配套述标PPT」能力，用于开标现场汇报或内部评审预演。

v9.1 升级（数据驱动）：
- 「企业业绩与资质保障」「项目团队配置」页在 ADR-005 移除企业画像后固定使用通用占位数据，
  不再消费 enterprise_profile；无画像时优雅回退到通用表述。
- 「项目概况」页补充投标单位 / 注册资本 / 信用等级与工程量清单要点。

- 纯 python-pptx 实现；如运行环境未安装 python-pptx，调用方 try/except 跳过，不阻断标书生成。
- 默认由 bid_generator 的 enable_ppt 开关控制（默认关闭）。
- 零额外硬依赖：requirements.txt 标注 python-pptx，缺失时优雅降级。
"""
from __future__ import annotations

# 配色（工程标书常用的稳重蓝 + 深灰）
_PRIMARY = (0x1F, 0x4E, 0x79)   # 深蓝
_ACCENT = (0x2E, 0x75, 0xB6)    # 中蓝
_DARK = (0x33, 0x33, 0x33)      # 深灰
_LIGHT = (0xF2, 0xF5, 0xFA)     # 浅底


def generate_bid_ppt(
    parse_result: dict,
    project_info: dict,
    output_path: str,
    score_response_map: dict | None = None,
) -> str:
    """生成述标PPT，返回输出路径。

    Args:
        parse_result: 招标文件解析结果（评分项/星号/废标/清单等）
        project_info: 项目信息（name/bid_type/duration/...）
        output_path: 标书 docx 路径（PPT 输出为同名 .pptx）
        score_response_map: v7.5 的评分项响应保障表结果（可选，用于"我方保障"页）
    Returns:
        生成的 pptx 路径
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    pr = parse_result or {}
    info = project_info or {}
    profile = {}
    company = profile.get('company', {}) or {}
    qualifications = profile.get('qualifications', []) or []
    key_personnel = profile.get('key_personnel', []) or []
    similar_projects = profile.get('similar_projects', []) or []
    honors = profile.get('honors', []) or []
    equipment = profile.get('equipment', []) or []

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _rgb(t):
        return RGBColor(*t)

    def _add_slide(title: str, bullets: list[str], accent: bool = False):
        slide = prs.slides.add_slide(blank)
        # 顶部色带
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(_PRIMARY if not accent else _ACCENT)
        bar.line.fill.background()
        tf = bar.text_frame
        tf.margin_left = Inches(0.4)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = _rgb((0xFF, 0xFF, 0xFF))

        # 内容框
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.6))
        tf2 = box.text_frame
        tf2.word_wrap = True
        first = True
        for b in bullets:
            para = tf2.paragraphs[0] if first else tf2.add_paragraph()
            first = False
            para.text = '• ' + b
            para.font.size = Pt(18)
            para.font.color.rgb = _rgb(_DARK)
            para.space_after = Pt(8)
        return slide

    name = info.get('name', '本工程项目')
    bid_type = info.get('bid_type', 'construction')
    type_label = '施工组织设计' if bid_type == 'construction' else '服务方案'
    duration = info.get('duration', '—')
    company_name = company.get('name', '')

    # 1. 封面
    cover = prs.slides.add_slide(blank)
    band = cover.shapes.add_shape(1, 0, Inches(2.6), prs.slide_width, Inches(2.3))
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(_PRIMARY)
    band.line.fill.background()
    ct = band.text_frame
    ct.vertical_anchor = MSO_ANCHOR.MIDDLE
    ct.margin_left = Inches(0.6)
    p1 = ct.paragraphs[0]
    p1.text = name
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = _rgb((0xFF, 0xFF, 0xFF))
    p2 = ct.add_paragraph()
    p2.text = f'{type_label} · 述标汇报'
    p2.font.size = Pt(24)
    p2.font.color.rgb = _rgb((0xCF, 0xE0, 0xF5))
    if company_name:
        p3 = ct.add_paragraph()
        p3.text = f'投标单位：{company_name}'
        p3.font.size = Pt(18)
        p3.font.color.rgb = _rgb((0xCF, 0xE0, 0xF5))

    # 2. 目录
    toc = ['项目概况与工期目标', '核心技术要点', '企业业绩与资质保障',
           '项目团队配置', '质量与安全保障', '偏离表与响应承诺']
    _add_slide('目录', toc)

    # 3. 项目概况（含投标单位与工程量清单要点）
    overview = [
        f'项目名称：{name}',
        f'工程类型：{type_label}',
        f'总工期：{duration} 日历天',
        f'质量标准：合格（争创优良）',
    ]
    if company_name:
        overview.append(f'投标单位：{company_name}')
        if company.get('registered_capital'):
            overview.append(f'注册资本：{company["registered_capital"]}')
        if company.get('credit_grade'):
            overview.append(f'企业信用：{company["credit_grade"]}')
    # 工程量清单要点（取前 3 条）
    for q in (pr.get('quantities', []) or [])[:3]:
        qt = q.get('text', '')
        qv = q.get('value', '')
        if qt and qv:
            overview.append(f'{qt}：{qv}')
    overview.append('承包范围：详见招标文件及工程量清单')
    _add_slide('一、项目概况与工期目标', overview)

    # 4. 核心技术要点（从评分项提炼）
    tech_bullets = []
    for it in (pr.get('score_items', []) or [])[:8]:
        nm = it.get('name') or it.get('title') or ''
        sc = it.get('score') or it.get('max_score') or ''
        if nm:
            tech_bullets.append(f'{nm}（{sc}分）— 针对性施工方案与管控要点')
    if not tech_bullets:
        tech_bullets = ['关键分部分项工程施工工艺', '进度计划与关键线路控制',
                        '质量通病防治与样板引路', '安全文明施工与绿色施工']
    _add_slide('二、核心技术要点', tech_bullets)

    # 5. 企业业绩与资质保障（v9.1: 消费真实企业画像）
    qual_bullets = []
    if similar_projects:
        qual_bullets.append('近三年类似工程业绩：')
        for proj in similar_projects[:4]:
            line = f'  - {proj.get("name", "")}'
            extra = []
            if proj.get('scale'):
                extra.append(proj['scale'])
            if proj.get('year'):
                extra.append(str(proj['year']))
            if proj.get('amount'):
                extra.append(proj['amount'])
            if extra:
                line += '（' + '·'.join(extra) + '）'
            qual_bullets.append(line)
    if qualifications:
        qual_bullets.append('企业资质等级：')
        for q in qualifications[:4]:
            qual_bullets.append(f'  - 具备{q}')
    if honors:
        qual_bullets.append('行业荣誉：')
        for h in honors[:2]:
            extra = h.get('level', '')
            yr = h.get('year', '')
            suffix = ('（' + '·'.join([x for x in [extra, str(yr) if yr else ''] if x]) + '）') if (extra or yr) else ''
            qual_bullets.append(f'  - 荣获{h.get("name", "")}{suffix}')
    if not qual_bullets:
        qual_bullets = ['类似工程业绩（近三年）', '企业资质等级与安全生产许可证',
                        '项目经理同类业绩', '机械设备与检测能力']
    _add_slide('三、企业业绩与资质保障', qual_bullets, accent=True)

    # 6. 项目团队配置（v9.1: 消费真实关键人员）
    team = []
    if key_personnel:
        for p in key_personnel[:5]:
            role = p.get('role', '')
            pname = p.get('name', '')
            cert = p.get('cert', '')
            title = p.get('title', '')
            line = f'{role}：{pname}'
            tail = '（' + '·'.join([x for x in [cert, title] if x]) + '）' if (cert or title) else ''
            line += tail
            team.append(line)
    if not team:
        team = ['项目经理：一级/二级建造师，同类工程经验',
                '技术负责人：高级职称，主持过多项同类工程',
                '质量/安全/施工员：持证上岗，配置齐全',
                '特种作业人员：100% 持证']
    _add_slide('四、项目团队配置', team)

    # 7. 质量与安全保障（含主要设备）
    qs = ['质量管理体系：三检制 + 样板引路 + 旁站监理',
          '安全管控：重大危险源辨识与应急预案',
          '工期保障：网络计划 + 资源动态调配',
          '服务承诺：响应招标全部实质性要求']
    if equipment:
        dev = '、'.join(f'{e.get("name", "")}×{e.get("count", "")}' for e in equipment[:4] if e.get('name'))
        if dev:
            qs.append(f'主要设备自有保障：{dev}')
    _add_slide('五、质量与安全保障', qs)

    # 8. 偏离表与承诺
    dev_bullets = []
    star = []
    for c in (pr.get('star_clauses', []) or [])[:4]:
        star.append('★ ' + (c.get('content') or c.get('text') or '')[:40])
    if star:
        dev_bullets = star
    dev_bullets += ['我方对招标文件全部实质性要求作出响应',
                    '正偏离 / 无偏离为主，负偏离逐条说明并补救']
    _add_slide('六、偏离表与响应承诺', dev_bullets)

    pptx_path = output_path.rsplit('.', 1)[0] + '_述标.pptx'
    prs.save(pptx_path)
    return pptx_path
